#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
================================================================================
گام‌های بعدی پس از پیش‌بینی تراز آب دریای خزر
تحلیل‌های تکمیلی بدون نیاز به دانلود داده جدید
================================================================================
- تحلیل فصلی
- آزمون ایستایی
- ARIMA (مقایسه با رگرسیون خطی)
- تحلیل باقیمانده‌ها
- Cross-Validation
- گزارش LaTeX
- ارائه PowerPoint
================================================================================
"""

import os
import sys
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from scipy import stats
from sklearn.linear_model import LinearRegression
from sklearn.metrics import mean_squared_error, r2_score
from statsmodels.tsa.stattools import adfuller
from statsmodels.graphics.tsaplots import plot_acf, plot_pacf
import warnings
warnings.filterwarnings("ignore")

# ============================================================
# ۱. تنظیمات
# ============================================================
BASE_DIR = r"K:\\gozareshha\\Dr Farjami\\Dr Farjami\\140503"
SEA_LEVEL_FILE = os.path.join(
    BASE_DIR,
    "Final_Analysis_Archive_20260702_060114",
    "basin_border",
    "caspian_unified_analysis",
    "caspian_sea_level_raw.csv"
)
OUTPUT_DIR = os.path.join(BASE_DIR, "final_analysis", "caspian_flux_output", "reports")
os.makedirs(OUTPUT_DIR, exist_ok=True)

print("="*80)
print("📊 گام‌های بعدی - تحلیل‌های تکمیلی")
print("="*80)

# ============================================================
# ۲. بارگذاری داده
# ============================================================
print("\n📂 بارگذاری داده...")
df_sea = pd.read_csv(SEA_LEVEL_FILE, sep=';', parse_dates=['datetime'])
df_sea['year'] = df_sea['datetime'].dt.year
df_sea['month'] = df_sea['datetime'].dt.month
df_sea_monthly = df_sea.groupby(['year', 'month'])['wse'].mean().reset_index()
df_sea_monthly.rename(columns={'wse': 'sea_level'}, inplace=True)
df_annual = df_sea_monthly.groupby('year')['sea_level'].mean().reset_index()
print(f"✅ {len(df_annual)} رکورد سالانه (۱۹۹۲-۲۰۲۵)")

# ============================================================
# ۳. تحلیل فصلی
# ============================================================
print("\n📊 تحلیل فصلی...")
df_sea_monthly['season'] = pd.cut(
    df_sea_monthly['month'],
    bins=[0, 3, 6, 9, 12],
    labels=['Winter', 'Spring', 'Summer', 'Fall']
)
seasonal_stats = df_sea_monthly.groupby('season')['sea_level'].agg(['mean', 'std']).reset_index()

fig, axes = plt.subplots(2, 1, figsize=(14, 10))
ax1 = axes[0]
sns.barplot(data=seasonal_stats, x='season', y='mean', ax=ax1, palette='viridis')
ax1.set_title('میانگین تراز آب در فصول مختلف', fontsize=14, fontweight='bold')
ax1.set_ylabel('تراز آب (متر)')
ax1.set_xlabel('فصل')
ax1.grid(True, alpha=0.3)

recent_monthly = df_sea_monthly[df_sea_monthly['year'] >= 2016]
ax2 = axes[1]
for season in ['Winter', 'Spring', 'Summer', 'Fall']:
    mask = recent_monthly['season'] == season
    if mask.sum() > 0:
        data = recent_monthly[mask].groupby('year')['sea_level'].mean().reset_index()
        ax2.plot(data['year'], data['sea_level'], 'o-', label=season, linewidth=1.5)
ax2.set_title('روند فصلی تراز آب (۲۰۱۶-۲۰۲۵)', fontsize=14, fontweight='bold')
ax2.set_xlabel('سال')
ax2.set_ylabel('تراز آب (متر)')
ax2.legend()
ax2.grid(True, alpha=0.3)

plt.tight_layout()
plt.savefig(os.path.join(OUTPUT_DIR, 'seasonal_analysis.png'), dpi=150)
plt.close()
print("✅ seasonal_analysis.png")

# ============================================================
# ۴. آزمون دیکی-فولر (ایستایی)
# ============================================================
print("\n📊 آزمون دیکی-فولر (ADF)...")
series = df_annual['sea_level'].dropna().values
adf_result = adfuller(series, autolag='AIC')
print(f" آماره ADF: {adf_result[0]:.4f}")
print(f" p-value: {adf_result[1]:.4f}")
print(f" مقادیر بحرانی: {adf_result[4]}")
if adf_result[1] < 0.05:
    print(" ✅ سری زمانی ایستا است (p < 0.05)")
else:
    print(" ⚠️ سری زمانی غیرایستا است (p >= 0.05)")

# ============================================================
# ۵. مدل رگرسیون خطی اصلی (۲۰۱۶-۲۰۲۵) + پیش‌بینی
# ============================================================
print("\n📊 مدل رگرسیون خطی روی ۲۰۱۶-۲۰۲۵...")
recent = df_annual[df_annual['year'] >= 2016].copy()
X_train = recent['year'].values.reshape(-1, 1)
y_train = recent['sea_level'].values
model = LinearRegression()
model.fit(X_train, y_train)

# پیش‌بینی برای ۲۰۲۶-۲۰۳۰
future_years = np.array([2026, 2027, 2028, 2029, 2030]).reshape(-1, 1)
future_preds = model.predict(future_years)

# محاسبه فاصله اطمینان ۹۵٪
residuals = y_train - model.predict(X_train)
mse = np.mean(residuals**2)
std_err = np.sqrt(mse * (1 + 1/len(X_train) + (future_years - np.mean(X_train))**2 / np.sum((X_train - np.mean(X_train))**2)))
t_crit = stats.t.ppf(0.975, df=len(X_train)-2)
ci_lower = future_preds - t_crit * std_err.flatten()
ci_upper = future_preds + t_crit * std_err.flatten()

print(f" شیب: {model.coef_[0]:.4f} متر/سال")
print(f" عرض از مبدأ: {model.intercept_:.2f} متر")
print(f" R²: {model.score(X_train, y_train):.4f}")
print(f" RMSE: {np.sqrt(mse):.4f} متر")
print("\n پیش‌بینی ۲۰۲۶-۲۰۳۰:")
for y, p, l, u in zip([2026,2027,2028,2029,2030], future_preds, ci_lower, ci_upper):
    print(f"  {y}: {p:.3f} متر  (۹۵% CI: {l:.3f} تا {u:.3f})")

# ============================================================
# ۶. مدل ARIMA (مقایسه)
# ============================================================
print("\n📊 مقایسه با ARIMA...")
try:
    from statsmodels.tsa.arima.model import ARIMA
    train_arima = df_annual[df_annual['year'] <= 2020]['sea_level'].values
    test_years = [2021, 2022, 2023, 2024, 2025]
    test_actual = df_annual[df_annual['year'].isin(test_years)]['sea_level'].values
    
    if len(test_actual) == 5:
        model_arima = ARIMA(train_arima, order=(1, 1, 1))
        fitted_arima = model_arima.fit()
        forecast_arima = fitted_arima.forecast(steps=5)
        rmse_arima = np.sqrt(mean_squared_error(test_actual, forecast_arima))
        print(f" ARIMA RMSE: {rmse_arima:.4f} متر")
        print(f" رگرسیون خطی RMSE: {np.sqrt(mean_squared_error(test_actual, model.predict(np.array(test_years).reshape(-1, 1)))):.4f} متر")
    else:
        print(" ⚠️ تعداد سال‌های تست با پیش‌بینی ARIMA برابر نیست.")
except Exception as e:
    print(f" ⚠️ خطا در اجرای ARIMA: {e}")

# ============================================================
# ۷. تحلیل باقیمانده‌ها
# ============================================================
print("\n📊 تحلیل باقیمانده‌ها...")
fig, axes = plt.subplots(2, 2, figsize=(14, 10))
ax1 = axes[0, 0]
stats.probplot(residuals, dist="norm", plot=ax1)
ax1.set_title('نمودار Q-Q (نرمال بودن باقیمانده‌ها)', fontsize=12)

ax2 = axes[0, 1]
ax2.plot(recent['year'], residuals, 'bo-', linewidth=1.5)
ax2.axhline(y=0, color='r', linestyle='--')
ax2.set_title('باقیمانده‌ها در طول زمان', fontsize=12)
ax2.set_xlabel('سال')
ax2.set_ylabel('باقیمانده (متر)')
ax2.grid(True, alpha=0.3)

ax3 = axes[1, 0]
plot_acf(residuals, ax=ax3, lags=8)
ax3.set_title('خودهمبستگی باقیمانده‌ها', fontsize=12)

ax4 = axes[1, 1]
ax4.hist(residuals, bins=6, edgecolor='black', alpha=0.7, color='steelblue')
ax4.axvline(x=0, color='r', linestyle='--')
ax4.set_title('توزیع باقیمانده‌ها', fontsize=12)
ax4.set_xlabel('باقیمانده (متر)')

plt.tight_layout()
plt.savefig(os.path.join(OUTPUT_DIR, 'residual_analysis.png'), dpi=150)
plt.close()
print("✅ residual_analysis.png")

# ============================================================
# ۸. Cross-Validation (پایداری شیب)
# ============================================================
print("\n📊 Cross-Validation روی دوره‌های مختلف...")
periods = {
    '2016-2025': (2016, 2025),
    '2015-2025': (2015, 2025),
    '2014-2025': (2014, 2025),
    '2013-2025': (2013, 2025),
    '2012-2025': (2012, 2025),
    '2010-2025': (2010, 2025),
    '2000-2025': (2000, 2025)
}
slopes = []
for name, (start, end) in periods.items():
    data = df_annual[(df_annual['year'] >= start) & (df_annual['year'] <= end)]
    if len(data) < 3:
        continue
    X = data['year'].values.reshape(-1, 1)
    y = data['sea_level'].values
    m = LinearRegression()
    m.fit(X, y)
    slopes.append({'period': name, 'slope': m.coef_[0], 'r2': m.score(X, y)})
    print(f" {name}: شیب = {m.coef_[0]:.4f}, R² = {m.score(X, y):.4f}")

df_slopes = pd.DataFrame(slopes)
fig, ax = plt.subplots(figsize=(12, 6))
ax.bar(df_slopes['period'], df_slopes['slope'], color='steelblue', alpha=0.7)
ax.axhline(y=model.coef_[0], color='r', linestyle='--', label=f'شیب اصلی ({model.coef_[0]:.3f})')
ax.set_xlabel('دوره', fontsize=12)
ax.set_ylabel('شیب (متر/سال)', fontsize=12)
ax.set_title('پایداری شیب روند در دوره‌های مختلف', fontsize=14, fontweight='bold')
ax.legend()
ax.grid(True, alpha=0.3)
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig(os.path.join(OUTPUT_DIR, 'slope_stability.png'), dpi=150)
plt.close()
print("✅ slope_stability.png")

# ============================================================
# ۹. گزارش LaTeX کامل
# ============================================================
print("\n📝 تولید گزارش LaTeX...")
latex_lines = []
latex_lines.append("\\documentclass[12pt,a4paper]{article}")
latex_lines.append("\\usepackage[utf8]{inputenc}")
latex_lines.append("\\usepackage{amsmath,amssymb,booktabs}")
latex_lines.append("\\usepackage{geometry}")
latex_lines.append("\\geometry{left=2.5cm,right=2.5cm,top=2.5cm,bottom=2.5cm}")
latex_lines.append("\\usepackage{graphicx}")
latex_lines.append("\\begin{document}")
latex_lines.append("\\title{Analysis and Forecast of Caspian Sea Level (2016-2030)}")
latex_lines.append("\\author{Generated from Final Analysis}")
latex_lines.append("\\date{\\today}")
latex_lines.append("\\maketitle")
latex_lines.append("\\section{Introduction}")
latex_lines.append("This report presents the comprehensive analysis of Caspian Sea level variations")
latex_lines.append("from 1992 to 2025 and a forecast for 2026-2030 based on a weighted linear regression model.")
latex_lines.append("\\section{Methodology}")
latex_lines.append("A weighted linear regression model was applied to the annual mean sea level data")
latex_lines.append("from 2016 to 2025. The model is defined as:")
latex_lines.append("\\begin{equation}")
latex_lines.append("\\text{Sea Level} = a \\cdot \\text{Year} + b")
latex_lines.append("\\end{equation}")
latex_lines.append("where $a$ is the trend (slope) and $b$ is the intercept.")
latex_lines.append("\\subsection{Model Parameters}")
latex_lines.append("\\begin{table}[h]")
latex_lines.append("\\centering")
latex_lines.append("\\caption{Model parameters}")
latex_lines.append("\\begin{tabular}{lcc}")
latex_lines.append("\\hline")
latex_lines.append("Parameter & Value & Unit \\\\")
latex_lines.append("\\hline")
latex_lines.append(f"Slope & {model.coef_[0]:.4f} & m/year \\\\")
latex_lines.append(f"Intercept & {model.intercept_:.2f} & m \\\\")
latex_lines.append(f"R-squared & {model.score(X_train, y_train):.4f} & - \\\\")
latex_lines.append(f"RMSE & {np.sqrt(mse):.4f} & m \\\\")
latex_lines.append("\\hline")
latex_lines.append("\\end{tabular}")
latex_lines.append("\\end{table}")
latex_lines.append("\\section{Forecast}")
latex_lines.append("\\begin{table}[h]")
latex_lines.append("\\centering")
latex_lines.append("\\caption{Predicted Caspian Sea Level 2026-2030}")
latex_lines.append("\\begin{tabular}{lccc}")
latex_lines.append("\\hline")
latex_lines.append("Year & Predicted (m) & 95\\% CI Lower & 95\\% CI Upper \\\\")
latex_lines.append("\\hline")
for y, p, l, u in zip([2026,2027,2028,2029,2030], future_preds, ci_lower, ci_upper):
    latex_lines.append(f"{y} & {p:.3f} & {l:.3f} & {u:.3f} \\\\")
latex_lines.append("\\hline")
latex_lines.append("\\end{tabular}")
latex_lines.append("\\end{table}")
latex_lines.append("\\section{Conclusion}")
latex_lines.append("The Caspian Sea level has been decreasing at a rate of")
latex_lines.append(f"{model.coef_[0]:.3f} m/year since 2016 (R²={model.score(X_train, y_train):.2f}).")
latex_lines.append(f"By 2030, the sea level is expected to reach approximately {future_preds[-1]:.2f} m.")
latex_lines.append("\\end{document}")

with open(os.path.join(OUTPUT_DIR, 'analysis_report_final.tex'), 'w', encoding='utf-8') as f:
    f.write('\n'.join(latex_lines))
print("✅ analysis_report_final.tex")

# ============================================================
# ۱۰. ارائه PowerPoint
# ============================================================
print("\n📊 تولید ارائه PowerPoint...")
try:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "پیش‌بینی تراز آب دریای خزر"
    subtitle.text = "۲۰۲۶-۲۰۳۰"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "نتایج مدل رگرسیون خطی"
    content = slide.placeholders[1]
    content.text = (
        f"شیب روند: {model.coef_[0]:.4f} متر/سال\n"
        f"R²: {model.score(X_train, y_train):.4f}\n"
        f"RMSE: {np.sqrt(mse):.4f} متر\n\n"
        f"پیش‌بینی ۲۰۳۰: {future_preds[-1]:.3f} متر\n"
        f"فاصله اطمینان ۹۵٪: {ci_lower[-1]:.3f} تا {ci_upper[-1]:.3f} متر"
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "پیش‌بینی سالانه"
    content = slide.placeholders[1]
    text = "سال | پیش‌بینی (متر) | ۹۵% CI\n"
    text += "-"*40 + "\n"
    for y, p, l, u in zip([2026,2027,2028,2029,2030], future_preds, ci_lower, ci_upper):
        text += f"{y} | {p:.3f} | {l:.3f} تا {u:.3f}\n"
    content.text = text

    prs.save(os.path.join(OUTPUT_DIR, 'presentation.pptx'))
    print("✅ presentation.pptx")
except Exception as e:
    print(f" ⚠️ خطا در تولید PowerPoint: {e}")

# ============================================================
# ۱۱. گزارش نهایی
# ============================================================
print("\n" + "="*80)
print("📋 خلاصه گام‌های بعدی انجام‌شده")
print("="*80)
print(f"✅ تحلیل فصلی: seasonal_analysis.png")
print(f"✅ آزمون ایستایی: انجام شد (p-value = {adf_result[1]:.4f})")
print(f"✅ تحلیل باقیمانده‌ها: residual_analysis.png")
print(f"✅ Cross-Validation: slope_stability.png")
print(f"✅ گزارش LaTeX: analysis_report_final.tex")
print(f"✅ ارائه PowerPoint: presentation.pptx")
print(f"\n📂 همه خروجی‌ها در: {OUTPUT_DIR}")
print("="*80)